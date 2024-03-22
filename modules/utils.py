import re
from typing import Any
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from scipy.ndimage import label ,center_of_mass
from matplotlib.axes._axes import Axes
import numpy.typing as npt
import numpy as np
import pickle
from matplotlib.colors import LinearSegmentedColormap
from sqlalchemy import create_engine
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px



def calculateExecutedTime(func):
    from time import time
    def wrapper(*args, **kwargs):
        s = time()
        result = func(*args, **kwargs)
        print(f"Executed Function Cost: {time() - s:.2f}s")
        return result
    return wrapper


def replace_str_time(str_time:str) -> str:
    """Using replace method to replace dash, colon and empty character to Null character in string of time.

    Params:
    -------
        str_time (str): string of time

    
    >>> str_time_format = '2023-06-25 15:48:45'
    >>> processed = replace_str_time(str_time_format)
    >>> processed
    20230625154845
    
    Returns:
        str: the processed string of time
    """

    if isinstance(str_time, str):
        return str_time.replace('-', '').replace(':', '').replace(' ', '')
    raise TypeError(f'The str_time format should be string not {type(str_time)}')
 


def regexes_str_time(str_time:str) -> str:
    """Using regexes method to replace dash, colon and empty character to Null character in string of time.

    Params:
    -------
        str_time (str): string of time

    
    >>> str_time_format = '2023-06-25 15:48:45'
    >>> processed = regexes_str_time(str_time_format)
    >>> processed
    20230625154845
    
    Returns:
        str: the processed string of time
    """
    if isinstance(str_time, str):
        return re.sub(r'-|:| |', '', str_time)
    else:
        raise TypeError(f'The str_time format should be string not {type(str_time)}')



def dict_to_df(dict_df:dict) -> pd.DataFrame:
    if len(dict_df) == 1:
        return pd.DataFrame(dict_df, index=[0])
    return pd.DataFrame(dict_df)



from pymongo import MongoClient
from pymongo.collection import Collection
from gridfs import GridFS
from bson import ObjectId

def connect_MongoDB(client:str, db_name:str, collection:str) -> Collection:
    client = MongoClient(client)
    user_db = client[db_name]
    user_collection = user_db[collection]
    return user_collection



def grid_fs(client:str, db_name:str, collection:str) -> GridFS:
    client = MongoClient(client)
    user_db = client[db_name]
    fs = GridFS(user_db, collection=collection)
    return fs


def get_df_from_period_of_time(start_date:str, end_date:str, InsType:str, OPID:str) -> pd.DataFrame:
    client = connect_MongoDB(client="mongodb://xxx:xxx@10.88.26.102:27017", db_name='MT', collection='LUM_SummaryTable')
    cursor = client.find(
        {
            'CreateTime': {'$gte':start_date, '$lte':end_date}, 
        },
        {
            '_id':0, 'CreateTime':1, 'SHEET_ID':1, 'OPID':1, 'LED_TYPE':1, 'NGCNT':1, 'LightingCheck_2D':1, 'Inspection_Type':1
        }
    )
    ts_df = pd.DataFrame.from_records(cursor).fillna('')
    ts_df = ts_df[(ts_df['Inspection_Type']==InsType) & (ts_df['OPID']==OPID)]
    ts_df = ts_df.sort_values(by=['CreateTime', 'SHEET_ID', 'NGCNT'], ascending=False)
    ts_df = ts_df.drop_duplicates(subset=['SHEET_ID', 'LED_TYPE', 'LightingCheck_2D'], keep='first').reset_index(drop=True)
    ts_df = ts_df.drop_duplicates(subset=['SHEET_ID', 'LED_TYPE'], keep='first').reset_index(drop=True)
    return ts_df




class plot_defect_info():
    """Plot the specific RGB defect info from specific Sheet ID."""
    def __init__(self, sheet_ID:str|None=None) -> None:
        self.sheet_ID = sheet_ID
        self.OrackeDB = "oracle://L4A_ULEDMFG_AP:L4A_ULEDMFG$AP@LTH41-VIP:1521/?service_name=L4AULEDH"
        self.MongoDBClient = "mongodb://xxx:xxx@10.88.26.102:27017"
        self.DB = "MT"
    
            
    def get_BONDING_MAP_df(self) -> pd.DataFrame:
        engine1 = create_engine(self.OrackeDB)
        conn1 = engine1.connect()
        if self.sheet_ID != None:
            sql = f'''
                    select sheet_id,
                    key_value_01 as coc2,
                    key_value_02 as area,
                    key_value_03 as wafer_id,
                    op_id,
                    lm_time
                    
                    from beolh.h_wip_sheetevent
                    where code_cat = 'BONDING_MAP'
                    and sheet_id = '{self.sheet_ID}'
                    order by lm_time
                    ''' 
                    
            bonding_map_df = pd.read_sql(sql, con=conn1) 
            conn1.close()

            bonding_map_df = bonding_map_df.drop_duplicates(['coc2', 'area'], keep='last').sort_values(
                ['coc2', 'lm_time'], ascending=True).reset_index(drop=True)
            
            return bonding_map_df
        
        
    def get_COC2_df(self) -> pd.DataFrame:
        bonding_map_df = self.get_BONDING_MAP_df()
        AOIclient = connect_MongoDB(client=self.MongoDBClient, db_name=self.DB, collection='COC2_AOI')
        # AOIfs= grid_fs(client='mongodb://xxx:xxx@10.88.26.102:27017', db_name='MT', collection='COC2_AOI_ARRAY')

        aoi_df_ls = []
        coc2_opid_ls = ['C2-ATO']
        
        for coc2 in tuple(dict.fromkeys(bonding_map_df['coc2'])):

            cursor = AOIclient.find(
                {
                    'SHEET_ID': coc2, 'OPID':{'$in': coc2_opid_ls}
                },
                {
                    '_id':0, 
                    'CreateTime':1, 
                    'SHEET_ID':1, 
                    'OPID':1, 
                    'LED_TYPE':1, 
                    'arr_id':1, 
                    'df_id':1, 
                    'CHIP':1
                }
            )
            aoi_temp_df = pd.DataFrame.from_records(cursor)
            aoi_df_ls.append(aoi_temp_df)

        coc2_df = pd.concat(aoi_df_ls)
        
        return coc2_df
    
    
    def get_COC2_df_without_TFT(self) -> pd.DataFrame:
        AOIclient = connect_MongoDB(client=self.MongoDBClient, db_name=self.DB, collection='COC2_AOI')
        coc2_opid_ls = ['C2-ATO', 'empty']
        
        cursor = AOIclient.find(
                {
                    'SHEET_ID': self.sheet_ID, 'OPID':{'$in': coc2_opid_ls}
                },
                {
                    '_id':0, 
                    'CreateTime':1, 
                    'SHEET_ID':1, 
                    'OPID':1, 
                    'LED_TYPE':1, 
                    'arr_id':1, 
                    'df_id':1, 
                    'CHIP':1
                }
            )
        
        coc2_df = pd.DataFrame.from_records(cursor)
        # print(coc2_df)
        
        return coc2_df

    
    def get_TFT_CreateTime_df(self) -> pd.DataFrame:
        TFTclient = connect_MongoDB(client=self.MongoDBClient, db_name=self.DB, collection='LUM_SummaryTable')

        ts = TFTclient.find(
            {'SHEET_ID':self.sheet_ID}, 
            {
                '_id':0, 
                'Grade':1, 
                'OPID':1, 
                'SHEET_ID':1,
                'CreateTime':1, 
                'NGCNT':1, 
                'LED_TYPE':1, 
                'Luminance_2D':1, 
                'ACTUAL_RECIPE':1,
                'LightingCheck_2D':1,
                'Yield':1, 
                'Inspection_Type':1, 
                'Dataframe_id':1
            }
        )
        
        ts_df = pd.DataFrame.from_records(ts)
        
        ts_df = ts_df.fillna('')
        ts_df = ts_df.sort_values(by=['CreateTime', 'LED_TYPE', 'NGCNT'], ascending=False).reset_index(drop=True)
        return ts_df
    
    
    def get_TFT_newest_df(self, df, InsType, OPID, recipe) -> pd.DataFrame:
        
        if recipe != "":
            ts_df = df[(df['Inspection_Type']==InsType) & (df['OPID']==OPID) & (df['ACTUAL_RECIPE']==recipe)]
        else:
            ts_df = df[(df['Inspection_Type']==InsType) & (df['OPID']==OPID)]
        
        ts_df = ts_df.drop_duplicates(subset=['SHEET_ID', 'LED_TYPE', 'LightingCheck_2D'], keep='first').reset_index(drop=True)
        ts_df = ts_df.drop_duplicates(subset=['SHEET_ID', 'LED_TYPE'], keep='first').reset_index(drop=True)
        return ts_df
    
    
    def get_option_list(self, ts_df:pd.DataFrame) -> list:
        """Let column of createTime and OPID into a list"""
        option_list = (ts_df['CreateTime']  + " " + ts_df['OPID']).unique()
        return option_list
    
       
        
    def get_specific_object_id(self, df:pd.DataFrame, column_name:str, LED_TYPE:str, Ins_tpye=str|None) -> ObjectId: 
        """Give a Dataframe and select the column what you want then give the LED_TYPE and Inspection type, 
        and you will get the specific ObjectID

        Args:
        ------
        
        df (pd.Dataframe): the dataframe including Object_ID
        column_name (str): the column you want to choose
        LED_TYPE: R,  G,  B
        Ins_type: L255,  L10, L0
        
        reture ObjectId
        """
        if Ins_tpye != None:
            object_id = df[(df['LED_TYPE']==LED_TYPE) & (df['Inspection_Type']==Ins_tpye)][column_name].tolist()[0]
            return object_id
        
        object_id = df[(df['LED_TYPE']==LED_TYPE)][column_name].tolist()[0]
        return object_id
    
    
    def get_defect_coord(self, object_id:str, fs:GridFS, coc2:False) -> tuple[list, list]:
        """
        Get the X-axis and Y-axis defect coordinates that will be used to plot defect scatter.

        Args:
            object_id (str): An ObjectID in MongoDB.
            fs (GridFS): The GridFS object.
            coc2 (bool): Determines if the array should be processed using specific steps. Defaults to False.

        Returns:
            tuple[list, list]: A tuple containing two lists - X-axis and Y-axis defect coordinates.
        """
        arr = fs.get(ObjectId(object_id)).read()
        arr = pickle.loads(arr)
        arr = np.flip(np.flip(arr, 0), 1)
        
        if coc2:
            arr = np.flip(np.flip(arr, 0), 1)
            y, x = arr.shape
            arr = np.flip(arr.reshape(x, y).T, 1)        
            arr = np.where(arr!=0, 1, arr) # defect code to 1
            arr = np.where(arr==1, 0, 1) # defect to 0
        
        arr = np.where(arr==0)
        x_coord = arr[1]
        y_coord = arr[0]
        
        return list(x_coord), list(y_coord)

    ### The production of the customization product has already ceased, so the corresponding function was deprecated in April 2023 ### 
    # def check_skip_pithes(self, arr:npt.ArrayLike, reduce_w:int, reduce_h:int) -> np.ndarray:
    #     """Change the array to dataframe, and check that whether skip pitches product.

    #     Params:
    #     -----------
    #         arr (npt.ArrayLike): R or G or B Yield 2-D array
    #         reduce_w (int): the width after skip pitches
    #         reduce_h (int): the height after skip pitches
    #     """
    #     if np.all(arr==1):
    #         return arr
        
    #     orgin_arr_df = pd.DataFrame(arr)
    #     temp_arr_df = orgin_arr_df.copy()
    #     skip_pitch_w = reduce_w
    #     skip_pitch_h = reduce_h
    #     cnt_w, cnt_h = 0, 0

    #     for i in range(len(orgin_arr_df.columns)):
    #         series_arr = np.asarray(orgin_arr_df[i], dtype=int)
    #         if np.all(series_arr==10) or np.all(series_arr==0):
    #             cnt_w += 1

    #     for _, row in orgin_arr_df.iterrows():
    #         row_arr = np.asarray(row, dtype=int)
    #         if np.all(row_arr==10) or np.all(row_arr==0):
    #             cnt_h += 1

    #     # 如果有跳 pitches 就將有跳過的地方改為1
    #     # 以原本的dataframe為基準, 將其改變後用 temp dataframe 取代
    #     if cnt_w == skip_pitch_w and cnt_h == skip_pitch_h:
    #         # 調整 column 的 serise
    #         for col in orgin_arr_df.columns:
    #             series_arr = np.asarray(temp_arr_df[col], dtype=int)
    #             if np.all(series_arr==10):
    #                 temp_arr_df[col] = np.where(temp_arr_df[col]==10, 1, temp_arr_df[col])
    #             elif np.all(series_arr==0):
    #                 temp_arr_df[col] = np.where(temp_arr_df[col]==0, 1, temp_arr_df[col])
    #         # 調整 rows 的 seriese
    #         for index, row in orgin_arr_df.iterrows():
    #             row_arr = np.asarray(row, dtype=int)
    #             if np.all(row_arr==10):
    #                 temp_arr_df.loc[index] = np.where(row_arr==10, 1, row_arr)
    #             elif np.all(row_arr==0):
    #                 temp_arr_df.loc[index] = np.where(row_arr==0, 1, row_arr)
    #         return np.asarray(temp_arr_df, dtype=int)
        
    #     return arr
    ### The production of the customization product has already ceased, so the corresponding function was deprecated in April 2023 ### 


    def get_heat_map_imshow_params(self, object_id:str, LED_TYPE:str, fs:GridFS) -> tuple[np.ndarray, LinearSegmentedColormap, int, int]:
        lum_arr = fs.get(ObjectId(object_id)).read()
        lum_arr = pickle.loads(lum_arr)
        
        if LED_TYPE =='R':
            colors = ['#000000', '#FF0000']
        elif LED_TYPE =='G':
            colors = ['#000000', '#00FF00']
        elif LED_TYPE =='B':
            colors = ['#000000', '#0000FF']

        cmap = LinearSegmentedColormap.from_list('CMAP', colors)
        lum_max = np.amax(lum_arr)
        lum_min = np.amin(lum_arr)
        
        return lum_arr, cmap, lum_max, lum_min


    def get_defect_imshow_params(self, object_id:str, LED_TYPE:str, fs:GridFS, coc2=False) -> tuple[np.ndarray, LinearSegmentedColormap]:
        defect_arr = fs.get(ObjectId(object_id)).read()
        defect_arr = pickle.loads(defect_arr)
        
        if coc2:
            y, x = defect_arr.shape
            defect_arr = np.flip(defect_arr.reshape(x, y).T, 1)        
            defect_arr = np.where(defect_arr!=0, 1, defect_arr) # defect code to 1
            defect_arr = np.where(defect_arr==1, 0, 1) # defect to 0
        
        NG_map = np.where(defect_arr==0, 1, 0)
        
        del defect_arr
        
        if LED_TYPE =='R':
            colors = ['red', 'black']
        elif LED_TYPE =='G':
            colors = ['green', 'black']
        elif LED_TYPE =='B':
            colors = ['blue', 'black']

        cmap = LinearSegmentedColormap.from_list('CMAP', colors) 
        
        return NG_map, cmap


    def get_NGCNT_Yield(self, df:pd.DataFrame, LED_TYPE:str) -> list:
        ins_ls = ['L255', 'edge_Dark_point', 'L0', 'L10']
        order_ls = []
        for ins in ins_ls:
            normal = df[(df['LED_TYPE']==LED_TYPE) & (df['Inspection_Type']==ins)]
            normal_ngcnt = normal.NGCNT.tolist()
            normal_yield = normal.Yield.tolist()
            final = normal_ngcnt + normal_yield
            order_ls.append(final)
            
        return order_ls


    def get_TFT_full_RGB_Dataframe(self, df:pd.DataFrame, Ins_tpye:str) -> tuple[pd.DataFrame, str, str]:
        sort_col_ls = ['LED_TYPE', 'Inspection_Type']
        duplicate_col = ['OPID', 'LED_TYPE', 'Inspection_Type']
        specific_time_df = df.drop_duplicates(duplicate_col, keep='first')
        specific_time_df = specific_time_df.sort_values(by=sort_col_ls, ascending=False).reset_index(drop=True)
        
        fs = grid_fs(self.MongoDBClient, self.DB, collection='LUM_SummaryTable')
        
        R_df_ID = self.get_specific_object_id(specific_time_df, 'Dataframe_id', 'R', Ins_tpye) 
        G_df_ID = self.get_specific_object_id(specific_time_df, 'Dataframe_id', 'G', Ins_tpye)
        B_df_ID = self.get_specific_object_id(specific_time_df, 'Dataframe_id', 'B', Ins_tpye)
        
        R_df = pickle.loads(fs.get(ObjectId(R_df_ID)).read())
        R_df = pd.DataFrame(R_df)
        
        G_df = pickle.loads(fs.get(ObjectId(G_df_ID)).read())
        G_df = pd.DataFrame(G_df)
        
        B_df = pickle.loads(fs.get(ObjectId(B_df_ID)).read())
        B_df = pd.DataFrame(B_df)
        
        full_df = pd.concat([R_df, G_df, B_df])
        
        specific_defect_code_col = ''
        specific_lumiance_col = ''
    
        if Ins_tpye == 'L255':
            specific_defect_code_col = 'Defect_Code'
            specific_lumiance_col = 'LED_Luminance'
            
        else:
            specific_defect_code_col = f'Defect_Code_{Ins_tpye}'
            specific_lumiance_col = f'{Ins_tpye}_LED_Luminance'
        
        return full_df, specific_defect_code_col, specific_lumiance_col
    
    
    def array_to_dataframe(self, Object_ID) -> pd.DataFrame:
        df = pickle.loads(Object_ID)
        if isinstance(df, np.ndarray):
            return pd.DataFrame(df)
        return df
    
    
    def seriesHaveChip(self, ChipColumn:pd.Series) -> bool:
        """The function will check whether the series includes "Chip" or not.

        Params:
        -----
            ChipColumn (pd.Series): The column of Chip

        Returns:
        -------
            bool
        """
        firstValue = ChipColumn.reset_index(drop=True)[0]
        if str(firstValue).startswith("Chip"):
            return True
        return False
    
    
    def concat_RGB_dataframe(self, coc2_df:pd.DataFrame) -> tuple[pd.DataFrame, int]:
        fs = grid_fs(self.MongoDBClient, self.DB, collection='COC2_AOI_DF')
        R_COC2_DF_ID = self.get_specific_object_id(coc2_df, 'df_id', 'R', Ins_tpye=None)
        G_COC2_DF_ID = self.get_specific_object_id(coc2_df, 'df_id', 'G', Ins_tpye=None)
        B_COC2_DF_ID = self.get_specific_object_id(coc2_df, 'df_id', 'B', Ins_tpye=None)
        
        R_coc2_df = fs.get(ObjectId(R_COC2_DF_ID)).read()
        R_coc2_df = self.array_to_dataframe(R_coc2_df)
        
        G_coc2_df = fs.get(ObjectId(G_COC2_DF_ID)).read()
        G_coc2_df = self.array_to_dataframe(G_coc2_df)
        
        B_coc2_df = fs.get(ObjectId(B_COC2_DF_ID)).read()
        B_coc2_df = self.array_to_dataframe(B_coc2_df)
        
        col_name = ['LED_TYPE', 'LED_Index_X', 'LED_Index_Y', 'Shift_X', 'Shif_Y' ,'Rotate', 'Defect_Reciepe', 'LINK', 'Chip']
        
        full_df = pd.concat([R_coc2_df, G_coc2_df, B_coc2_df])
        max_index_x = full_df[7].max()
        full_df = full_df[full_df[13] != 'OK']
        
        full_df.drop(columns=[0, 1, 2, 3, 4, 5, 12, 14], inplace=True)
        
        full_df = full_df[full_df.columns[:9]]
        full_df.columns = col_name
        
        return full_df, max_index_x
        
    
    def get_COC2_full_RGB_Dataframe(self, onlyCOC2=True) -> pd.DataFrame:
        """The function will get dataframe including RGB's data, if that including Chip data
        then merge it.

        Returns:
            pd.DataFrame: The dataframe including RGB's data or Chip
        """
        if onlyCOC2:
            coc2_df = self.get_COC2_df_without_TFT()
        else:
            coc2_df = self.get_COC2_df()
        
        if 'CHIP' in coc2_df.columns:
            haveChip = self.seriesHaveChip(coc2_df['CHIP'])
            if haveChip:
                chip_df_ls = []
                chip_items = coc2_df['CHIP'].sort_values(ascending=True).unique()
                for chip in chip_items:
                    chip_df, _ = self.concat_RGB_dataframe(coc2_df=coc2_df[coc2_df['CHIP']==chip])
                    chip_df_ls.append(chip_df)
                full_df = pd.concat(chip_df_ls)
                
                return full_df
            return self.concat_RGB_dataframe(coc2_df=coc2_df)
        return self.concat_RGB_dataframe(coc2_df=coc2_df)
    
    
    def interact_scatter(
            self, 
            df, 
            col_led_index_x, 
            col_led_index_y, 
            col_of_color, 
            symbol: Any|None=None, 
            hover_data:list|None=None
        ) -> go.Figure:
        if 'LED_Index_J' in df.columns:
            max_y = df['LED_Index_J'].max()
        else:
            max_y = df['LED_Index_Y'].max()
        
        if symbol != None:
            df = df[df[symbol]!= '']
        
        color_discrete_map = {'R': 'rgb(255,0,0)', 'G': 'rgb(0,255,0)', 'B': 'rgb(0,0,255)'}
        fig = px.scatter(
            df,
            x=col_led_index_x,
            y=col_led_index_y,
            color=col_of_color,
            color_discrete_map=color_discrete_map,
            symbol=symbol,
            hover_data=hover_data,
            color_continuous_scale="reds",
            width=1000,
            height=500,
            range_y=[max_y, 0],
        )
        
        fig.update_layout(
            xaxis={'side': 'top'},
        )
        
        return fig


class Mark_cluster():
    """
    Mark_cluster
    =============

    The function use frame to mark the cluster defect from 2D array, and the threshold of defect count will be 
    show on the frame base on input value.
    
    
    Params
    ------
    
    defect_map (npt.ArrayLike): An binary 2-D array
    threshold_point (int): An number use to choose extract the defect in a row
    axs (Axes): matplotlib.axes._axes
    
    How to use the class::
    ---------------------
    
    Create a 2-D array::
    
      >>> fig, ax = plt.subplots(2,2)
      >>> array = np.ones((540, 240))

    Simulate defect distribution::
    
      >>> array[1:20, 50:200] = 0 
      >>> mark_cluster = Mark_cluster(array, threshold, ax[0,0])
      >>> ax[0, 0] = mark_cluster()
      >>> ax[0, 0].imshow(array)
    
    """
    def __init__(self, defect_map:npt.ArrayLike, threshold_point:int, axs:Axes) -> None:
        self.defect_map = defect_map
        self.threshold_point = threshold_point,
        self.axs = axs
   
   
    def __call__(self, defect_map:npt.ArrayLike, threshold_point:int, axs:Axes) -> Axes:
        self.defect_map = defect_map
        self.threshold_point = threshold_point,
        self.axs = axs
        self.axs = self.get_marked_axes()
        
        return self.axs
    
    
    def get_marked_axes(self) -> Axes:
        self.defect_map = np.asarray(self.defect_map, dtype='uint8')
        structure = np.ones((3,3), dtype=bool)
        labeled_array, num_features = label(self.defect_map, structure=structure)
    
        for i in range(num_features):
            
            # calculate cluster point in total
            size = np.sum(labeled_array == i+1)

            if size >= self.threshold_point:
                # find centers_of_cluster index
                center = center_of_mass(self.defect_map, labeled_array, i+1)
                center_y, center_x = int(center[0]), int(center[1])

                # calculate frame size
                y, x = np.where(labeled_array == i+1)

                rect_height = (np.max(y) - np.min(y)) + 10
                rect_width = (np.max(x) - np.min(x)) + 10

                # plot the retangle of frame
                self.axs.add_patch(plt.Rectangle(
                        (center_x - rect_width/2, center_y - rect_height/2), 
                        rect_width, 
                        rect_height, 
                        fill=False, 
                        edgecolor='lime', 
                        lw=2, 
                    )
                )
                
                # show the defect count upper the frame
                self.axs.text((center_x-rect_width/2)+3, (center_y-rect_height/2)-3, size, color='black', fontsize=8, bbox=dict(facecolor='lime', pad=2, edgecolor='none'))
                
        return self.axs

    
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from glob import glob
import os
from io import BytesIO

class CreatePPT():
    """
    Create PPT file by image name
    """
    def __init__(self):
        self.template_ppt_path = './modules/template ppt/template.pptx'
        self.prs = Presentation(self.template_ppt_path)
        self.imgls = glob('./modules/temp/*.png')
        self.temp_ppt_path = "./modules/temp/temp.pptx"
    
    
    def __call__(self, image_path) -> Any:
        
        if os.path.exists(self.temp_ppt_path):
            self.prs = Presentation(self.temp_ppt_path)
        else:
            self.prs = Presentation(self.template_ppt_path)
            
        self.addSlide(image_path)
        self.save()
    
        
    def addSlide(self, image_path:str, left=Inches(0), top=Inches(1.5), width=None, height=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_picture(image_path, left=left, top=top, width=width, height=height)
    
    
    def save(self):
        self.prs.save("./modules/temp/temp.pptx")
    
    
    def delete_temp_file(self):
        if os.path.exists(self.temp_ppt_path):
            os.remove(self.temp_ppt_path)

    
    def load_ppt(self):
        if os.path.exists(self.temp_ppt_path):
            document = Presentation(self.temp_ppt_path)
        else:
            document = Presentation(self.template_ppt_path)
            
        binary_output = BytesIO()
        document.save(binary_output)
        
        return binary_output
      
            
from scipy import signal
class conv_cluster():
    def __init__(self, arr: np.ndarray, kernel_size: int, threshold: int) -> Any:
        self.arr = arr
        self.kernel_size = kernel_size
        self.threshold = threshold
        
        
    def __call__(self, arr: np.ndarray, kernel_size: int, threshold: int) -> Any:
        self.arr = arr
        self.kernel_size = kernel_size
        self.threshold = threshold
    
    
    def create_kernel(self) -> np.ndarray:
        kernel = np.ones((self.kernel_size, self.kernel_size), dtype=np.uint8)
        
        return kernel
    
    
    # def TraditionConvolution(self, padding=0, stride=1) -> np.ndarray:
    #     img_height, img_width = self.arr.shape
    #     kernel = self.create_kernel()
    #     kernel_y, kernel_x = kernel.shape
        
    #     if kernel_y > img_height or kernel_x > img_width:
    #         raise ValueError("Error: The size of the kernel matrix exceeds the size of the image")
        
    #     # stride = m - n + 1 
    #     # m: image size, n: kernel size
    #     stride_x = img_width - kernel_x + 1
    #     stride_y = img_height - kernel_y + 1
        
    #     m = int((arrx - kx + (2*padding)) / stride) + 1
    #     n = int((arry - ky + (2*padding)) / stride) + 1

    #     newImage = np.zeros((n, m))
        
    #     for i in range(stride_y):
    #         for j in range(stride_x):
    #             newImage[i, j] = np.sum(arr[i:i+kernel_y, j:j+kernel_x])
                
    #     return newImage
                
    
    def convImage(self) -> np.ndarray:
        kernel = self.create_kernel()            
        img = signal.convolve2d(self.arr, kernel, mode='valid')
        return img


    def checkOutofSpec(self, convImage: np.ndarray) -> int:
        res = np.any(convImage >= self.threshold)
        if res == True:
            return 1
        return 0

        
    def plotConvClusterFrame(self, sheetID: str, convImage: np.ndarray, pixelType: str|None=""):
        """Plot conv cluster

        Args:
            sheetID (str): product ID
            convImage (np.ndarray): 2-D array
            pixelType (str | None, optional): Main pixel or by pixel. Defaults to None.

        Returns:
            Figure
        """
        fig, axs = plt.subplots(1,1)
        fig.set_figheight(5)
        fig.set_figwidth(17.5)
        fig.suptitle(f'{sheetID} Conv Cluster Result {pixelType}', fontsize=16, y=1)
        
        colors = ['white', 'black']
        cmap = mcolors.LinearSegmentedColormap.from_list('CMAP', colors)
        axs.imshow(self.arr, cmap=cmap, aspect='auto')
        
        axs.set_title(f'Ksize:{self.kernel_size}*{self.kernel_size} Th:{self.threshold}', y=-0.1)
        axs.tick_params(top=True, labeltop=True, bottom=False, labelbottom=False)
        
        for i in range(convImage.shape[0]):
            for j in range(convImage.shape[1]):
                if convImage[i, j] >= self.threshold:
                    coordx = i-0.5
                    coordy = j-0.5
                    axs.add_patch(plt.Rectangle(
                            ((coordy, coordx)), 
                            self.kernel_size, 
                            self.kernel_size, 
                            fill=False, 
                            edgecolor='lime', 
                            lw=2, 
                        )
                    )

                    # show the defect count upper the frame
                    axs.text(
                        coordy, 
                        coordx, 
                        int(convImage[i, j]), 
                        color='black', 
                        fontsize=10, 
                        bbox=dict(facecolor='lime', pad=2, edgecolor='none')
                    )
                    
                else:
                    continue
                
        return fig

               
class CornerOffDot:
    def __call__(self, matrix, step):
        self.matrix = matrix
        self.step = step
        
        y, x = matrix.shape
        
        # 取左上角的3*3
        top_left_corner = matrix[0:step, 0:step]

        # 取右上角的3*3
        top_right_corner = matrix[0:step, x-step:x]

        # 取左下角的3*3
        bottom_left_corner = matrix[y-step:y, 0:step]

        # 取右下角的3*3
        bottom_right_corner = matrix[y-step:y, x-step:x]

        # extent=[x0, x1, y1, y0]
        self.corner_range_dict = {
            '[0][0]' : ['top_left_corner', top_left_corner, [1, step+1, step+1, 1]],
            '[0][1]' : ['top_right_corner', top_right_corner, [x-step, x, step+1, 1]],
            '[1][0]' : ['bottom_left_corner', bottom_left_corner, [1, step+1, y, y-step]],
            '[1][1]' : ['bottom_right_corner', bottom_right_corner, [x-step, x, y, y-step]]
        }
        
        fig = self.plotCorner()
        
        
    def plotCorner(self):
        colors = ['white', 'black']
        cmap = mcolors.LinearSegmentedColormap.from_list('CMAP', colors)
        fig, axs = plt.subplots(2, 2)
        # fig.set_figheight(8)
        # fig.set_figwidth(8)
        for i in range(2):
            for j in range(2):
                params_list = self.corner_range_dict.get(f'[{i}][{j}]')
                axs[i, j].set_title(params_list[0], y=-0.2)
                axs[i, j].imshow(params_list[1], cmap=cmap, extent=params_list[2])
                axs[i, j].tick_params(top=True, labeltop=True, bottom=False, labelbottom=False)

        plt.tight_layout()
        
        return fig   
    
    
    
    
          



